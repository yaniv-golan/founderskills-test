#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx>=1.0"]
# ///
"""Read a PowerPoint deck's text when it cannot be rendered to PDF.

This is the DEGRADED path and exists only as a fallback. A `.pptx` is binary, so the
Read tool cannot open it at all; the good path converts to PDF with LibreOffice and
reviews the rendered slides, which is what a reader of the deck actually sees. Use this
only when no converter is available.

What it recovers: slide text, table cells, and SPEAKER NOTES — which matter more than
their obscurity suggests, because notes routinely carry the narrative the slide only
gestures at.

What it cannot recover, and why the caller must act on that: anything about how the deck
LOOKS. Layout, density, hierarchy, and any figure that lives inside an image are all
invisible here. A review built on this output must set `input_format` to `"text"` so the
Design & Readability criteria gate to `not_applicable` rather than being scored by a
reviewer that never saw a slide.

Emits JSON on stdout; `-o` writes it to a file as well.
"""

from __future__ import annotations

import argparse
import json
import pathlib
import sys
from typing import Any


def _shape_text(shape: Any) -> list[str]:
    out: list[str] = []
    if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
        out.extend(ln for ln in shape.text_frame.text.splitlines() if ln.strip())
    if getattr(shape, "has_table", False) and shape.has_table:
        for row in shape.table.rows:
            cells = [c.text.replace("\n", " ").strip() for c in row.cells]
            if any(cells):
                out.append(" | ".join(cells))
    # A NATIVE chart carries its own numbers, and they are otherwise lost: a live run over
    # a real deck detected three doughnut charts, read their captions, and never pulled the
    # values behind them. Read them here, paired to their category labels so nothing has to
    # align two lists by position.
    if getattr(shape, "has_chart", False) and shape.has_chart:
        chart = shape.chart
        try:
            cats = [str(c) for c in chart.plots[0].categories]
        except (IndexError, ValueError):
            cats = []
        # Whether the numbers are PRINTED on the slide is a different question from whether
        # they exist in the file, and the caller needs both: a plotted-but-unlabelled value
        # is real, but it is not something the deck can be said to state.
        try:
            shown = bool(chart.plots[0].data_labels.show_value)
        except (IndexError, AttributeError, ValueError):
            shown = False
        kind = str(chart.chart_type).split(" ")[0]
        out.append(f"[chart: {kind}; values {'printed on the slide' if shown else 'not printed'}]")
        for ser in chart.series:
            vals = ["" if v is None else f"{v:g}" for v in ser.values]
            name = (ser.name or "").strip() or "(unnamed series)"
            if cats and len(cats) == len(vals):
                pairs = ", ".join(f"{c}={v}" for c, v in zip(cats, vals, strict=True) if v != "")
                out.append(f"series {name}: {pairs}")
            elif any(vals):
                out.append(f"series {name}: " + " | ".join(v for v in vals if v))
    return out


def extract(path: pathlib.Path) -> dict[str, Any]:
    from pptx import Presentation

    prs = Presentation(str(path))
    slides: list[dict[str, Any]] = []
    images = 0
    for n, slide in enumerate(prs.slides, 1):
        lines: list[str] = []
        for shape in slide.shapes:
            lines.extend(_shape_text(shape))
            if shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
                images += 1
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append({"number": n, "text": "\n".join(lines), "speaker_notes": notes})
    return {
        "source": path.name,
        "total_slides": len(slides),
        # Surfaced so the caller can tell the founder what was NOT read. A deck carrying
        # many images is exactly the case where text-only review is weakest, and staying
        # quiet about it would present a partial review as a complete one.
        "images_not_read": images,
        "slides": slides,
    }


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("pptx", type=pathlib.Path)
    ap.add_argument("-o", "--output", type=pathlib.Path)
    ap.add_argument("--pretty", action="store_true")
    args = ap.parse_args()

    if not args.pptx.is_file():
        print(f"error: no such file: {args.pptx}", file=sys.stderr)
        return 1
    try:
        data = extract(args.pptx)
    except ImportError:
        print(
            "error: python-pptx is not installed, so this deck cannot be read as text either. "
            "Ask the founder to re-export the deck as PDF.",
            file=sys.stderr,
        )
        return 2
    except Exception as exc:  # noqa: BLE001 - a malformed upload must not look like a clean empty deck
        print(f"error: could not read {args.pptx.name}: {exc}", file=sys.stderr)
        return 1

    text = json.dumps(data, indent=2 if args.pretty else None)
    if args.output:
        args.output.write_text(text)
        print(json.dumps({"written": str(args.output), "total_slides": data["total_slides"]}))
    else:
        print(text)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
